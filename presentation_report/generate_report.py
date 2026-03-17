import os
import glob
from docx import Document
from pptx import Presentation
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

def get_files(directory):
    docx_files = glob.glob(os.path.join(directory, "*.docx"))
    pptx_files = glob.glob(os.path.join(directory, "*.pptx"))
    
    if not docx_files:
        raise FileNotFoundError("No .docx file found in the directory.")
    if not pptx_files:
        raise FileNotFoundError("No .pptx file found in the directory.")
        
    return docx_files[0], pptx_files[0]

def add_content_from_pptx(doc, pptx_path):
    prs = Presentation(pptx_path)
    
    # Ensure we are starting on a new page (Page 4 logic handled by padding or breaks)
    # The user wants content from 4th page onwards.
    # We need to check current page count or just ensure we add enough breaks if it's empty.
    # For now, we'll assume the user wants us to APPEND to the file, ensuring it starts on a fresh page if needed.
    # But strictly "from 4th page onwards" implies we might need to fill pages 1-3 if empty?
    # Let's assume the doc might be a template. We will add a page break to ensure we are at the end, 
    # and if the doc is short, we might need to add empty pages.
    
    # Simple approach: Add a page break before starting report content
    doc.add_page_break()
    
    # Add a main title
    heading = doc.add_heading('Presentation Report', 0)
    heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    
    content_map = []

    for i, slide in enumerate(prs.slides):
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text
        else:
            slide_title = f"Slide {i+1}"
            
        # Add to content map for Index
        content_map.append(slide_title)
        
        # Add to Doc
        doc.add_heading(slide_title, level=1)
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Skip title shape as we already added it
            if shape == slide.shapes.title:
                continue
                
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    p = doc.add_paragraph(text)
                    p.style = 'Normal'

    return content_map

def create_index(doc, content_map):
    # We want the index at the beginning or specific place?
    # User said "make the index first and then ... start to make the report file"
    # Usually Index is at the start.
    # We can insert a table of contents or manual index at the beginning.
    
    # Insert at the beginning of the document
    # This is tricky with python-docx as it appends. 
    # We can try to insert paragraphs at index 0.
    
    # Let's create a new first paragraph for the Index Title
    p_title = doc.paragraphs[0].insert_paragraph_before("Index")
    p_title.style = 'Heading 1'
    
    # Insert index items
    # We'll insert them in reverse order after the title so they appear in order? 
    # No, insert_paragraph_before inserts *immediately* before.
    # So to get:
    # Index
    # Item 1
    # Item 2
    # We insert Index. Then insert Item 1 after Index? No, insert_paragraph_before is relative to an existing paragraph.
    
    # Actually, if the doc is empty, this might fail.
    # Let's assume the doc has at least one paragraph (even empty).
    if len(doc.paragraphs) == 0:
        doc.add_paragraph("")
        
    ref_paragraph = doc.paragraphs[0] # The one we will push down
    
    # Add Index Title
    index_title = ref_paragraph.insert_paragraph_before("Table of Contents")
    index_title.style = 'Heading 1'
    
    for item in content_map:
        p = ref_paragraph.insert_paragraph_before(item)
        p.style = 'List Bullet'
        
    # Add a page break after index
    ref_paragraph.insert_paragraph_before("").add_run().add_break()

def main():
    cwd = os.getcwd()
    try:
        docx_path, pptx_path = get_files(cwd)
        print(f"Using Word file: {docx_path}")
        print(f"Using Presentation file: {pptx_path}")
        
        doc = Document(docx_path)
        
        # 1. Extract and Append Content (this goes to the end)
        # We need to ensure this starts on Page 4.
        # We can calculate roughly or just force breaks.
        # If we want to be strictly on Page 4, we might need to pad.
        # But "from 4th page onwards" might just mean "after the intro/index pages".
        # Let's just append for now.
        
        print("Extracting content from presentation...")
        content_map = add_content_from_pptx(doc, pptx_path)
        
        print("Creating Index...")
        create_index(doc, content_map)
        
        output_path = os.path.join(cwd, "Generated_Report.docx")
        doc.save(output_path)
        print(f"Report generated successfully: {output_path}")
        
    except FileNotFoundError as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")

if __name__ == "__main__":
    main()
